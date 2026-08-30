import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from typing import Optional

def generate_pptx_report(
    filename: str,
    target_col: str,
    kpis: dict,
    narrative_text: str,
    profile: dict,
    anomaly_metrics: dict = {},
    forecast_metrics: dict = {},
    segmentation_metrics: dict = {},
) -> Optional[bytes]:
    try:
        prs = Presentation()
        
        # SLIDE 1: Title
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Análisis Ejecutivo: {filename}"
        subtitle.text = f"MIO - Inteligencia Artificial\nGenerado el {datetime.now().strftime('%d/%m/%Y')}"
        
        # Style Title
        title.text_frame.paragraphs[0].font.name = 'Arial'
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(129, 90, 225) # mio-violet
        
        # SLIDE 2: KPIs & Profile
        bullet_slide_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_slide_layout)
        shapes2 = slide2.shapes
        title2 = shapes2.title
        body2 = shapes2.placeholders[1]
        
        title2.text = "Métricas Principales"
        tf2 = body2.text_frame
        
        n_rows = profile.get("n_rows", "N/A")
        n_cols = profile.get("n_cols", "N/A")
        tf2.text = f"Dataset: {n_rows} registros, {n_cols} columnas."
        
        p = tf2.add_paragraph()
        p.text = f"Columna Objetivo: {target_col}"
        p.level = 1
        
        for k, v in kpis.items():
            p = tf2.add_paragraph()
            p.text = f"{k}: {v}"
            p.level = 1
            
        p2 = tf2.add_paragraph()
        p2.text = f"Calidad de Datos: {profile.get('quality_score', 100)}/100 ({profile.get('quality_label', 'Alta')})"
        p2.level = 1

        # SLIDE 3: AI Narrative
        slide3 = prs.slides.add_slide(bullet_slide_layout)
        shapes3 = slide3.shapes
        title3 = shapes3.title
        body3 = shapes3.placeholders[1]
        
        title3.text = "Insights Generados por IA"
        # Narrative text is usually long, so we just drop it into the text frame
        clean_narrative = (
            narrative_text.replace("**", "")
            .replace("##", "")
            .replace("#", "")
        )
        body3.text_frame.text = clean_narrative
        
        # SLIDE 4: Machine Learning Highlights
        if anomaly_metrics or forecast_metrics:
            slide4 = prs.slides.add_slide(bullet_slide_layout)
            shapes4 = slide4.shapes
            title4 = shapes4.title
            body4 = shapes4.placeholders[1]
            
            title4.text = "Resultados Predictivos (Machine Learning)"
            tf4 = body4.text_frame
            tf4.text = "Resumen de Modelos"
            
            if forecast_metrics and "tendencia_pct" in forecast_metrics:
                p = tf4.add_paragraph()
                p.text = f"Proyección: Tendencia esperada del {forecast_metrics['tendencia_pct']}% en {forecast_metrics.get('periodos', 60)} días."
                p.level = 1
                
            if anomaly_metrics and "n_anomalias" in anomaly_metrics:
                p = tf4.add_paragraph()
                p.text = f"Anomalías: {anomaly_metrics['n_anomalias']} registros atípicos detectados ({anomaly_metrics.get('pct_anomalias', 0)}%)."
                p.level = 1
                
        # Save to bytes
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream.read()
        
    except Exception as e:
        print(f"Error generando PPTX: {e}")
        return None
